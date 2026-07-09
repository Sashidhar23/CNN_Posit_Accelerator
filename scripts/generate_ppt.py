"""
Generate a simple progress PPT for the CNN Posit Accelerator project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── helpers ──────────────────────────────────────────────────────────────

BG_COLOR   = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy
ACCENT     = RGBColor(0x00, 0xD2, 0xFF)    # cyan accent
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
ORANGE     = RGBColor(0xFF, 0xA5, 0x00)
GREEN      = RGBColor(0x00, 0xE6, 0x76)
YELLOW     = RGBColor(0xFF, 0xD7, 0x00)
DARK_BOX   = RGBColor(0x22, 0x22, 0x3A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, bullet_color=ACCENT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)

        # bullet character
        run_b = p.add_run()
        run_b.text = "▸ "
        run_b.font.size = Pt(font_size)
        run_b.font.color.rgb = bullet_color
        run_b.font.name = font_name

        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = color
        run_t.font.name = font_name
    return tf


def add_code_box(slide, left, top, width, height, code_lines, font_size=11):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x0D, 0x0D, 0x1A)
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)

    for i, line in enumerate(code_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(1)
        p.space_before = Pt(1)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Consolas"
        # color keywords
        if line.strip().startswith("//"):
            run.font.color.rgb = RGBColor(0x6A, 0x99, 0x55)
        elif any(kw in line for kw in ["module ", "parameter ", "input ", "output ",
                                        "wire ", "reg ", "assign ", "always ",
                                        "generate", "for ", "begin", "end",
                                        "posit_", "if ", "else"]):
            run.font.color.rgb = RGBColor(0x56, 0x9C, 0xD6)
        else:
            run.font.color.rgb = RGBColor(0xCE, 0xD4, 0xDA)
    return tf


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(0.04))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


def add_box_shape(slide, left, top, width, height, text, fill_color=DARK_BOX,
                  text_color=WHITE, font_size=14, bold=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = ACCENT
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.bold = bold
    run.font.name = "Calibri"

    # vertical centering
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    return shape


def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape

def add_down_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                                   Inches(left), Inches(top),
                                   Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Previous Work: Model Training & Weight Extraction
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(slide)

add_text_box(slide, 0.6, 0.3, 12, 0.7,
             "Previous Work — Model Training & Weight Extraction",
             font_size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.6, 0.95, 5)

# Left column — what was done
add_bullet_list(slide, 0.6, 1.3, 6, 4.5, [
    "Trained VGG16 on MNIST dataset (~31 MB weights)",
    "Trained GoogLeNet on Fashion-MNIST (~40 MB weights)",
    "Both models saved as PyTorch .pth state dictionaries",
    "Weights include convolution kernels, biases, batch-norm parameters",
    "Python reference model created for Posit<16,1> encoding/decoding",
    "float_to_posit16_1()  converts FP32 → 16-bit Posit (ES=1)",
], font_size=17)

# Right column — code snippet
add_text_box(slide, 7.2, 1.3, 5.5, 0.5,
             "posit_encoder.py  (Python reference)",
             font_size=14, color=YELLOW, bold=True)

add_code_box(slide, 7.2, 1.8, 5.5, 4.5, [
    "def float_to_posit16_1(x):",
    "    # 1. Handle sign",
    "    sign = 1 if x < 0 else 0",
    "",
    "    # 2. Compute regime k",
    "    #    USEED = 2^(2^ES) = 4",
    "    k = 0",
    "    while x >= USEED:  # x >= 4",
    "        x /= USEED;  k += 1",
    "",
    "    # 3. Extract exponent bit",
    "    exponent = 1 if x >= 2 else 0",
    "",
    "    # 4. Extract fraction bits",
    "    fraction = x - 1",
    "",
    "    # 5. Pack: sign | regime | exp | frac",
    "    return posit_bits",
])


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Previous Work: Weight Conversion Pipeline
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.6, 0.3, 12, 0.7,
             "Previous Work — FP32 to Posit Weight Conversion",
             font_size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.6, 0.95, 5)

# Flow diagram: .pth → Script → .npy
add_box_shape(slide, 1.0, 1.8, 2.5, 1.0,
              "VGG16 / GoogLeNet\n.pth weights\n(FP32)", font_size=14, bold=True)
add_arrow(slide, 3.7, 2.05, 0.8, 0.4)
add_box_shape(slide, 4.7, 1.8, 3.2, 1.0,
              "convert_convnext_weights.py\nnp.vectorize(float_to_posit16_1)",
              font_size=13, text_color=YELLOW)
add_arrow(slide, 8.1, 2.05, 0.8, 0.4)
add_box_shape(slide, 9.1, 1.8, 2.8, 1.0,
              "Posit<16,1> Weights\n.npy arrays\n(16-bit integers)", font_size=14, bold=True,
              text_color=GREEN)

# Explanation bullets
add_bullet_list(slide, 0.6, 3.5, 11.5, 3.5, [
    "Each FP32 weight/bias is individually converted to a 16-bit Posit<16,1> integer",
    "Conversion uses the Python reference encoder (float_to_posit16_1)",
    "Converted weights are stored in NumPy .npy format, preserving original tensor shapes",
    "check_converted_weights.py verifies conversion accuracy by decoding back to float",
    "These Posit-encoded weights are what the hardware accelerator will load into PEs",
], font_size=17)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Posit Decoder & Encoder (Verilog)
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.6, 0.3, 12, 0.7,
             "Verilog — Posit Decoder & Encoder",
             font_size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.6, 0.95, 4)

# Decoder section
add_text_box(slide, 0.6, 1.2, 5.5, 0.5,
             "posit_decoder.v — Unpacks N-bit Posit into fields",
             font_size=16, color=YELLOW, bold=True)

add_code_box(slide, 0.6, 1.7, 5.8, 3.2, [
    "module posit_decoder #(N=8, ES=1) (",
    "  input  [N-1:0] posit_in,",
    "  output sign, is_zero, is_nar,",
    "  output signed [$clog2(N):0] k,    // regime value",
    "  output [ES-1:0] exponent,",
    "  output [N-1:0]  fraction,",
    "  output [$clog2(N):0] frac_len",
    ");",
    "",
    "// Steps:",
    "// 1. Check special cases (zero, NaR)",
    "// 2. Extract sign, 2's complement for negative",
    "// 3. Count regime run-length → k",
    "// 4. Extract exponent bits after regime",
    "// 5. Extract fraction bits (left-justified)",
])

# Encoder section
add_text_box(slide, 6.8, 1.2, 6, 0.5,
             "posit_encoder.v — Packs fields back into N-bit Posit",
             font_size=16, color=YELLOW, bold=True)

add_code_box(slide, 6.8, 1.7, 5.8, 3.2, [
    "module posit_encoder #(N=8, ES=1) (",
    "  input  sign, is_zero, is_nar,",
    "  input  signed [$clog2(N):0] k,",
    "  input  [ES-1:0] exponent,",
    "  input  [N-1:0]  fraction,",
    "  output [N-1:0]  posit_out",
    ");",
    "",
    "// Steps:",
    "// 1. Handle zero / NaR special outputs",
    "// 2. Build regime bits from k",
    "//    k>=0 → (k+1) ones + terminating 0",
    "//    k<0  → |k| zeros + terminating 1",
    "// 3. Insert exponent, fraction bits",
    "// 4. Apply sign via 2's complement",
])

# Visual: bit layout
add_text_box(slide, 0.6, 5.2, 12, 0.5,
             "Posit<8,1> Bit Layout Example:",
             font_size=16, color=WHITE, bold=True)

# Draw bit fields as boxes
labels = [("S", 0.7), ("Regime", 2.2), ("E", 0.7), ("Fraction", 2.8)]
x = 0.8
for label, w in labels:
    col = ORANGE if label == "S" else (GREEN if label == "Regime" else (YELLOW if label == "E" else ACCENT))
    add_box_shape(slide, x, 5.7, w, 0.55, label, fill_color=DARK_BOX, text_color=col, font_size=16, bold=True)
    x += w + 0.15

add_text_box(slide, 7.5, 5.75, 5, 0.5,
             "sign | 1...10 (run) | exp bit(s) | mantissa bits",
             font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — Posit Adder & Multiplier
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.6, 0.3, 12, 0.7,
             "Verilog — Posit Adder & Multiplier",
             font_size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.6, 0.95, 4)

# Adder
add_text_box(slide, 0.6, 1.2, 6, 0.5,
             "posit_adder.v  (Combinational, ~300 lines)",
             font_size=16, color=YELLOW, bold=True)

add_bullet_list(slide, 0.6, 1.7, 5.8, 3.5, [
    "Decodes both inputs via posit_decoder",
    "Computes unified scale:  scale = k × 2^ES + exponent",
    "Builds significands with hidden 1:  sig = {1, fraction}",
    "Aligns smaller operand by right-shifting",
    "Adds/subtracts signed significands",
    "Normalizes result (find leading-one, shift)",
    "Decomposes result scale → new k and exponent",
    "Encodes output via posit_encoder",
], font_size=16)

# Multiplier
add_text_box(slide, 6.8, 1.2, 6, 0.5,
             "posit_multiplier.v  (Combinational, ~325 lines)",
             font_size=16, color=YELLOW, bold=True)

add_bullet_list(slide, 6.8, 1.7, 5.8, 3.5, [
    "Decodes both inputs via posit_decoder",
    "Result sign = sign_a XOR sign_b",
    "Builds mantissas:  mant = {1'b1, fraction}",
    "Integer multiply:  product = mant_a × mant_b",
    "Scales add:  total_scale = scale_a + scale_b",
    "Normalizes: if product MSB set → shift right, scale+1",
    "Trims trailing zeros for frac_len",
    "Saturates k to valid Posit range, encodes output",
], font_size=16)

# Common pattern box
add_text_box(slide, 0.6, 5.5, 12, 0.5,
             "Common Pattern:  Decode → Operate on (sign, scale, significand) → Normalize → Encode",
             font_size=18, color=GREEN, bold=True)

add_text_box(slide, 0.6, 6.0, 12, 0.6,
             "Both are fully combinational (no clk). Each adder/multiplier internally instantiates its own decoder and encoder.",
             font_size=15, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Processing Element (PE)
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.6, 0.3, 12, 0.7,
             "Verilog — Processing Element (PE)",
             font_size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.6, 0.95, 4)

add_text_box(slide, 0.6, 1.1, 6, 0.5,
             "pe.v — Weight-Stationary, 2-Stage Pipelined MAC",
             font_size=16, color=YELLOW, bold=True)

# Left: code
add_code_box(slide, 0.6, 1.6, 6.0, 4.3, [
    "module pe #(N=8, ES=1) (",
    "  input  clk, reset, pe_en, clear_acc, wshift,",
    "  input  [N-1:0] input_in, weight_in, psum_in,",
    "  output [N-1:0] input_out, psum_out, pe_output",
    ");",
    "",
    "  // Weight stays frozen (only loads when wshift=1)",
    "  always @(posedge clk)",
    "    if (wshift) weight_reg <= weight_in;",
    "",
    "  // Activation slides through (left → right)",
    "  always @(posedge clk)",
    "    if (pe_en) input_reg <= input_in;",
    "  assign input_out = input_reg;",
    "",
    "  // Stage 1: Multiply (combinational)",
    "  posit_multiplier MUL(.a(input_reg),.b(weight_reg),",
    "                       .out(product_comb));",
    "",
    "  // Stage 2: Add psum (combinational)",
    "  posit_adder ADD(.a(product_reg), .b(psum_pipe_reg),",
    "                  .out(psum_comb));",
    "",
    "  // Pipeline registers between stages",
    "  always @(posedge clk) begin",
    "    product_reg   <= product_comb; // stage 1 result",
    "    psum_pipe_reg <= psum_in;      // matching psum",
    "    psum_out_reg  <= psum_comb;    // stage 2 result",
    "  end",
], font_size=11)

# Right: diagram
add_text_box(slide, 7.2, 1.3, 5.5, 0.5,
             "PE Dataflow Diagram",
             font_size=16, color=YELLOW, bold=True)

# Input from left
add_box_shape(slide, 7.3, 2.0, 2.0, 0.6, "Activation In\n(from left PE)", font_size=12, text_color=ACCENT)
add_arrow(slide, 9.4, 2.1, 0.5, 0.3)
add_box_shape(slide, 10.0, 2.0, 1.5, 0.6, "input_reg", font_size=13, text_color=GREEN, bold=True)
add_arrow(slide, 11.6, 2.1, 0.5, 0.3)
add_text_box(slide, 12.1, 2.0, 1.0, 0.5, "→ Right", font_size=12, color=ACCENT)

# Weight
add_box_shape(slide, 10.0, 2.85, 1.5, 0.6, "weight_reg\n(stationary)", font_size=12, text_color=ORANGE, bold=True)

# Multiplier
add_down_arrow(slide, 10.55, 3.5, 0.35, 0.35)
add_box_shape(slide, 9.5, 3.95, 2.5, 0.55, "posit_multiplier", font_size=13, text_color=YELLOW, bold=True)

# Pipeline reg
add_down_arrow(slide, 10.55, 4.55, 0.35, 0.3)
add_box_shape(slide, 9.5, 4.95, 2.5, 0.45, "product_reg (pipeline)", font_size=12, text_color=LIGHT_GRAY, bold=True)

# Adder
add_down_arrow(slide, 10.55, 5.45, 0.35, 0.3)
# psum_in arrow
add_box_shape(slide, 7.3, 5.15, 1.8, 0.5, "psum_in\n(from above)", font_size=11, text_color=ACCENT)
add_arrow(slide, 9.2, 5.25, 0.3, 0.2)

add_box_shape(slide, 9.5, 5.85, 2.5, 0.55, "posit_adder", font_size=13, text_color=YELLOW, bold=True)

# Output
add_down_arrow(slide, 10.55, 6.45, 0.35, 0.3)
add_text_box(slide, 10.0, 6.8, 2.0, 0.5, "psum_out ↓", font_size=13, color=GREEN, bold=True)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Systolic Array
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.6, 0.3, 12, 0.7,
             "Verilog — 6×6 Weight-Stationary Systolic Array",
             font_size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.6, 0.95, 5)

# Left: description
add_bullet_list(slide, 0.6, 1.3, 5.5, 3.5, [
    "systolic_array.v  generates a ROWS×COLS grid of PEs",
    "Activations flow horizontally (left → right) via act_bus",
    "Partial sums flow vertically (top → bottom) via psum_bus",
    "Top row psum_in initialized to zero",
    "Bottom row psum_out = final column results",
    "Weights loaded in parallel via wshift before computation",
    "All PEs share clk, reset, pe_en, clear_acc signals",
], font_size=16)

# Left: systolic_core
add_text_box(slide, 0.6, 4.5, 5.5, 0.5,
             "systolic_core.v wraps the array with:",
             font_size=16, color=YELLOW, bold=True)

add_bullet_list(slide, 0.6, 5.0, 5.5, 2.0, [
    "Input FIFOs — per-row buffers for activation skewing",
    "Row delay registers — stagger activations diagonally",
    "Output delay registers — de-skew the column results",
    "This ensures correct data alignment across the pipeline",
], font_size=15, bullet_color=GREEN)

# Right: PE grid visual (simplified 4x4 representation)
add_text_box(slide, 7.0, 1.3, 5, 0.5,
             "Array Layout (6×6 PE Grid):",
             font_size=16, color=YELLOW, bold=True)

grid_rows = 4
grid_cols = 4
pe_w = 1.1
pe_h = 0.7
start_x = 8.0
start_y = 2.2

for row in range(grid_rows):
    # Activation arrow label
    add_text_box(slide, start_x - 1.0, start_y + row * (pe_h + 0.25) - 0.05,
                 0.9, 0.4, f"Act[{row}]→", font_size=11, color=ACCENT, bold=True)
    for col in range(grid_cols):
        x = start_x + col * (pe_w + 0.15)
        y = start_y + row * (pe_h + 0.25)
        label = f"PE\n[{row},{col}]"
        add_box_shape(slide, x, y, pe_w, pe_h, label, font_size=10,
                      text_color=WHITE, bold=True)

# Down arrows for psum labels below grid
for col in range(grid_cols):
    x = start_x + col * (pe_w + 0.15) + pe_w/2 - 0.25
    y = start_y + grid_rows * (pe_h + 0.25) - 0.1
    add_text_box(slide, x, y, 0.6, 0.4, "↓", font_size=16, color=GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER)

add_text_box(slide, start_x, start_y + grid_rows * (pe_h + 0.25) + 0.2, 5, 0.4,
             "psum_out  (final results per column)",
             font_size=13, color=GREEN, bold=True)

# PE quire variant note
add_text_box(slide, 7.0, 6.2, 5.5, 0.8,
             "Quire variants (pe_quire.v, systolic_array_quire.v)\n"
             "accumulate in a 128-bit fixed-point quire register\n"
             "for zero rounding error during MAC accumulation.",
             font_size=14, color=LIGHT_GRAY)


# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Summary & Status
# ═══════════════════════════════════════════════════════════════════════

slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)

add_text_box(slide, 0.6, 0.3, 12, 0.7,
             "Current Status & Next Steps",
             font_size=30, color=ACCENT, bold=True)
add_accent_line(slide, 0.6, 0.95, 4)

# Completed
add_text_box(slide, 0.6, 1.3, 6, 0.5,
             "✅  Completed", font_size=20, color=GREEN, bold=True)

add_bullet_list(slide, 0.6, 1.8, 5.8, 3.5, [
    "Posit<16,1> Python reference encoder & decoder",
    "Model training (VGG16, GoogLeNet) and weight extraction",
    "FP32 → Posit weight conversion pipeline",
    "Verilog posit_decoder — bit unpacking",
    "Verilog posit_encoder — bit packing",
    "Verilog posit_adder — full posit addition",
    "Verilog posit_multiplier — full posit multiplication",
    "Verilog PE — 2-stage pipelined weight-stationary PE",
    "Verilog PE Quire — PE with 128-bit quire accumulation",
    "Verilog 6×6 Systolic Array (standard + quire variants)",
    "Verilog Systolic Core with input FIFOs and skewing logic",
    "Testbenches written for all modules",
], font_size=15, bullet_color=GREEN)

# Next steps
add_text_box(slide, 7.0, 1.3, 5.5, 0.5,
             "🔲  Next Steps", font_size=20, color=ORANGE, bold=True)

add_bullet_list(slide, 7.0, 1.8, 5.5, 4.5, [
    "Run & verify all testbenches (PE, Systolic Array)",
    "Analyze simulation waveforms for timing correctness",
    "Build top-level controller / FSM for the accelerator",
    "Implement layer-level dataflow (tiling, feature maps)",
    "Load converted Posit weights into systolic array",
    "End-to-end inference test with MNIST/FMNIST images",
    "FPGA synthesis and resource utilization analysis",
    "Performance benchmarking vs FP32 / INT8 baselines",
], font_size=15, bullet_color=ORANGE)


# ── save ─────────────────────────────────────────────────────────────

import os
out_path = os.path.join(
    os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
    "outputs",
    "progress_presentation.pptx"
)
os.makedirs(os.path.dirname(out_path), exist_ok=True)
prs.save(out_path)
print(f"Saved to: {out_path}")
